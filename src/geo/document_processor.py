"""
文档处理模块 - 支持解析 Word、PDF、PPT、TXT 等格式
"""

import os
import re
import json
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime

# 尝试导入文档解析库
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from src.db.models import get_connection
from src.config import get_document_storage_dir


class DocumentProcessor:
    """文档处理器"""

    # 支持的文件类型
    SUPPORTED_TYPES = {
        '.txt': 'text',
        '.md': 'markdown',
        '.docx': 'word',
        '.doc': 'word',
        '.pdf': 'pdf',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint'
    }

    # 分类/标签选项（现在是自由输入，保留一些常用的供参考）
    SUGGESTED_TAGS = ['价格', '周期', '技术', '专利', '临床', '人群', '品牌', '案例', '数据']

    def __init__(self):
        pass

    def get_supported_extensions(self) -> List[str]:
        """获取支持的文件扩展名列表"""
        return list(self.SUPPORTED_TYPES.keys())

    def detect_file_type(self, filename: str) -> Optional[str]:
        """检测文件类型"""
        ext = Path(filename).suffix.lower()
        return self.SUPPORTED_TYPES.get(ext)

    def parse_document(self, file_path: str, file_type: str) -> Tuple[str, int]:
        """
        解析文档，返回 (文本内容, 字数)
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        if file_type == 'text' or file_type == 'markdown':
            return self._parse_txt(path)
        elif file_type == 'word':
            return self._parse_docx(path)
        elif file_type == 'pdf':
            return self._parse_pdf(path)
        elif file_type == 'powerpoint':
            return self._parse_pptx(path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

    def _parse_txt(self, path: Path) -> Tuple[str, int]:
        """解析TXT/Markdown文件"""
        try:
            content = path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            content = path.read_text(encoding='gbk', errors='ignore')
        return content, len(content)

    def _parse_docx(self, path: Path) -> Tuple[str, int]:
        """解析Word文档"""
        if not HAS_DOCX:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(path)
        paragraphs = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # 也提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    paragraphs.append(' | '.join(row_cells))

        content = '\n'.join(paragraphs)
        return content, len(content)

    def _parse_pdf(self, path: Path) -> Tuple[str, int]:
        """解析PDF文件"""
        if not HAS_PDF:
            raise ImportError("请安装 PyPDF2: pip install PyPDF2")

        content_parts = []
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    text = page.extract_text()
                    if text:
                        content_parts.append(text)
                except Exception:
                    continue

        content = '\n'.join(content_parts)
        return content, len(content)

    def _parse_pptx(self, path: Path) -> Tuple[str, int]:
        """解析PPT文件"""
        if not HAS_PPTX:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(path)
        content_parts = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            slide_content.append(f"--- 第 {slide_idx + 1} 页 ---\n")

            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_content.append(shape.text)

                # 提取表格
                if hasattr(shape, 'table'):
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_cells:
                            slide_content.append(' | '.join(row_cells))

            if slide_content:
                content_parts.append('\n'.join(slide_content))

        content = '\n\n'.join(content_parts)
        return content, len(content)

    def split_into_chunks(self, content: str, chunk_size: int = 500, overlap: int = 100) -> List[str]:
        """
        将长文本切分成片段（保留语义完整性）
        """
        # 先按段落分割
        paragraphs = content.split('\n')
        chunks = []
        current_chunk = []
        current_length = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # 如果当前段落很长，单独切分
            if len(para) > chunk_size:
                # 先保存当前chunk
                if current_chunk:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0

                # 切分长段落（按句子）
                sentences = re.split(r'([。！？.!?])', para)
                long_chunk_parts = []
                long_length = 0

                for i in range(0, len(sentences), 2):
                    sentence = sentences[i]
                    if i + 1 < len(sentences):
                        sentence += sentences[i + 1]

                    if long_length + len(sentence) > chunk_size and long_chunk_parts:
                        chunks.append(''.join(long_chunk_parts))
                        long_chunk_parts = [sentence]
                        long_length = len(sentence)
                    else:
                        long_chunk_parts.append(sentence)
                        long_length += len(sentence)

                if long_chunk_parts:
                    chunks.append(''.join(long_chunk_parts))

            elif current_length + len(para) > chunk_size and current_chunk:
                chunks.append('\n'.join(current_chunk))
                current_chunk = [para]
                current_length = len(para)
            else:
                current_chunk.append(para)
                current_length += len(para)

        if current_chunk:
            chunks.append('\n'.join(current_chunk))

        return chunks

    # ========== 数据库操作 ==========

    def save_document(self, project_id: Optional[int], original_filename: str,
                     storage_path: str, file_type: str, file_size: int,
                     tags: str = '') -> int:
        """保存文档记录（tags 用逗号分隔）"""
        conn = get_connection()
        cursor = conn.cursor()

        # 生成一个安全的文件名
        safe_filename = f"{int(datetime.now().timestamp())}_{Path(original_filename).name}"

        try:
            cursor.execute('''
                INSERT INTO geo_documents
                (project_id, filename, original_filename, file_type, file_size, category, storage_path, is_parsed)
                VALUES (?, ?, ?, ?, ?, ?, ?, 0)
            ''', (project_id, safe_filename, original_filename, file_type, file_size, tags, storage_path))

            doc_id = cursor.lastrowid
            conn.commit()
            return doc_id
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def update_document_parsed(self, doc_id: int, content: str, word_count: int):
        """更新文档解析状态"""
        conn = get_connection()
        cursor = conn.cursor()

        # 取前200字作为预览
        preview = content[:200] + '...' if len(content) > 200 else content

        try:
            cursor.execute('''
                UPDATE geo_documents
                SET content_preview = ?, word_count = ?, is_parsed = 1, updated_at = CURRENT_TIMESTAMP
                WHERE id = ?
            ''', (preview, word_count, doc_id))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def save_chunks(self, doc_id: int, chunks: List[str]):
        """保存文档片段"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            for idx, chunk in enumerate(chunks):
                cursor.execute('''
                    INSERT INTO geo_document_chunks
                    (document_id, chunk_index, content, content_length, is_embedded)
                    VALUES (?, ?, ?, ?, 0)
                ''', (doc_id, idx, chunk, len(chunk)))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def get_documents(self, project_id: Optional[int], tag: Optional[str] = None) -> List[Dict]:
        """获取文档列表（tag 可选，模糊匹配）"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            query = 'SELECT * FROM geo_documents WHERE project_id IS ?'
            params = [project_id]

            if tag:
                query += ' AND category LIKE ?'
                params.append(f'%{tag}%')

            query += ' ORDER BY created_at DESC'
            cursor.execute(query, params)
            docs = [dict(row) for row in cursor.fetchall()]
            # 把 category 字段改名为 tags
            for doc in docs:
                doc['tags'] = doc.get('category', '')
            return docs
        finally:
            conn.close()

    def get_document(self, doc_id: int) -> Optional[Dict]:
        """获取单个文档"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('SELECT * FROM geo_documents WHERE id = ?', (doc_id,))
            row = cursor.fetchone()
            return dict(row) if row else None
        finally:
            conn.close()

    def get_document_chunks(self, doc_id: int) -> List[Dict]:
        """获取文档的所有片段"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('''
                SELECT * FROM geo_document_chunks
                WHERE document_id = ? ORDER BY chunk_index
            ''', (doc_id,))
            return [dict(row) for row in cursor.fetchall()]
        finally:
            conn.close()

    def delete_document(self, doc_id: int):
        """删除文档（级联删除片段和摘要）"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('DELETE FROM geo_document_chunks WHERE document_id = ?', (doc_id,))
            cursor.execute('DELETE FROM geo_summaries WHERE summary_level = ? AND target_id = ?', ('document', doc_id))
            cursor.execute('DELETE FROM geo_documents WHERE id = ?', (doc_id,))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    # ========== 摘要操作 ==========

    def save_summary(self, project_id: Optional[int], summary_level: str,
                    target_id: Optional[int], title: str, content: str,
                    is_manual_edit: bool = False) -> int:
        """保存摘要"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('''
                INSERT INTO geo_summaries
                (project_id, summary_level, target_id, title, content, is_manual_edit)
                VALUES (?, ?, ?, ?, ?, ?)
            ''', (project_id, summary_level, target_id, title, content, 1 if is_manual_edit else 0))
            summary_id = cursor.lastrowid
            conn.commit()
            return summary_id
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def update_summary(self, summary_id: int, content: str, is_manual_edit: bool = False):
        """更新摘要"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('''
                UPDATE geo_summaries
                SET content = ?, is_manual_edit = ?, updated_at = CURRENT_TIMESTAMP
                WHERE id = ?
            ''', (content, 1 if is_manual_edit else 0, summary_id))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def get_summaries(self, project_id: Optional[int], summary_level: Optional[str] = None) -> List[Dict]:
        """获取摘要列表"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            query = 'SELECT * FROM geo_summaries WHERE project_id IS ?'
            params = [project_id]

            if summary_level:
                query += ' AND summary_level = ?'
                params.append(summary_level)

            query += ' ORDER BY created_at DESC'
            cursor.execute(query, params)
            return [dict(row) for row in cursor.fetchall()]
        finally:
            conn.close()

    def delete_summary(self, summary_id: int):
        """删除摘要"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('DELETE FROM geo_summaries WHERE id = ?', (summary_id,))
            conn.commit()
        except Exception as e:
            conn.rollback()
            raise e
        finally:
            conn.close()

    def get_all_chunks_for_project(self, project_id: Optional[int]) -> List[Dict]:
        """获取项目的所有文档片段（用于检索）"""
        conn = get_connection()
        cursor = conn.cursor()

        try:
            cursor.execute('''
                SELECT c.*, d.filename, d.original_filename, d.category
                FROM geo_document_chunks c
                JOIN geo_documents d ON c.document_id = d.id
                WHERE d.project_id IS ?
                ORDER BY d.created_at DESC, c.chunk_index
            ''', (project_id,))
            return [dict(row) for row in cursor.fetchall()]
        finally:
            conn.close()
